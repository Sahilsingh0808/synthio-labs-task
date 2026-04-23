"""Extract text from an uploaded presentation artifact.

Supported:
* PDF  — via ``pypdf``
* PPTX — via ``python-pptx``
* TXT / MD — decoded as UTF-8
* PNG / JPG / WEBP — passed as an image to the multimodal chat endpoint

For images we round-trip through the chat model so the caller receives plain
text (which is then fed into :func:`slide_generator.generate_deck`). This keeps
the deck-generation pipeline provider-agnostic.
"""

from __future__ import annotations

import base64
import io
import logging
from typing import Tuple

from azure_client import chat_deployment, get_client, token_kwargs

logger = logging.getLogger(__name__)

MAX_CHARS = 20_000


def _from_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    chunks = []
    for page in reader.pages:
        try:
            chunks.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n\n".join(c for c in chunks if c.strip())


def _from_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    chunks = []
    for i, slide in enumerate(prs.slides):
        pieces = [f"Slide {i + 1}:"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        pieces.append(f"- {text}")
        if len(pieces) > 1:
            chunks.append("\n".join(pieces))
    return "\n\n".join(chunks)


def _from_image(data: bytes, content_type: str) -> str:
    b64 = base64.b64encode(data).decode("ascii")
    client = get_client()
    model = chat_deployment()

    messages = [
        {
            "role": "system",
            "content": (
                "You are an expert at reading slides. Extract the title, bullet "
                "points, and any visible body text from the image. Preserve order. "
                "Return plain text only, no markdown."
            ),
        },
        {
            "role": "user",
            "content": [
                {"type": "text", "text": "Extract the slide content from this image."},
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:{content_type};base64,{b64}"},
                },
            ],
        },
    ]

    resp = client.chat.completions.create(
        model=model,
        messages=messages,
        **token_kwargs(1200, model),
    )
    return (resp.choices[0].message.content or "").strip()


def extract_text(filename: str, content_type: str, data: bytes) -> Tuple[str, str]:
    """Return ``(kind, text)`` where ``kind`` is a short label for logging."""

    name = (filename or "").lower()
    ctype = (content_type or "").lower()

    if name.endswith(".pdf") or "pdf" in ctype:
        return "pdf", _from_pdf(data)[:MAX_CHARS]

    if name.endswith(".pptx") or "presentation" in ctype:
        return "pptx", _from_pptx(data)[:MAX_CHARS]

    if name.endswith((".txt", ".md")) or ctype.startswith("text/"):
        return "text", data.decode("utf-8", errors="ignore")[:MAX_CHARS]

    if ctype.startswith("image/") or name.endswith((".png", ".jpg", ".jpeg", ".webp")):
        return "image", _from_image(data, ctype or "image/png")[:MAX_CHARS]

    raise ValueError(f"Unsupported file type: {filename} ({content_type})")
